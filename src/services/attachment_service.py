"""AttachmentService - Comprehensive attachment handling for EWS MCP v3.0."""

import logging
import io
from typing import Optional
from pathlib import Path

from ..core.attachment import Attachment, AttachmentContent, AttachmentType
from ..utils import safe_get


class AttachmentService:
    """
    Service for handling all attachment formats.

    Supports:
    - PDF (with pdfplumber)
    - DOCX (with python-docx)
    - Excel (with openpyxl)
    - PPTX (with python-pptx)
    - ZIP archives
    - Plain text
    - HTML
    """

    def __init__(self, ews_client):
        """
        Initialize AttachmentService.

        Args:
            ews_client: EWSClient instance
        """
        self.ews_client = ews_client
        self.logger = logging.getLogger(__name__)

    async def read_attachment(
        self,
        message_id: str,
        attachment_name: str,
        extract_text: bool = True,
        extract_images: bool = False,
        extract_metadata: bool = True,
        max_pages: int = 50
    ) -> Optional[AttachmentContent]:
        """
        Read and parse attachment content.

        Args:
            message_id: Parent message ID
            attachment_name: Name of attachment to read
            extract_text: Extract text content
            extract_images: Extract images
            extract_metadata: Extract metadata
            max_pages: Maximum pages for PDF

        Returns:
            AttachmentContent object or None
        """
        self.logger.info(f"Reading attachment: {attachment_name} from message {message_id}")

        try:
            # Get message
            message = await self._get_message(message_id)
            if not message:
                self.logger.error(f"Message not found: {message_id}")
                return None

            # Find attachment
            attachments = safe_get(message, 'attachments', []) or []
            target_attachment = None

            for att in attachments:
                if safe_get(att, 'name') == attachment_name:
                    target_attachment = att
                    break

            if not target_attachment:
                self.logger.error(f"Attachment not found: {attachment_name}")
                return None

            # Create Attachment metadata
            attachment = Attachment.from_ews_attachment(target_attachment, message_id)

            # Get content
            content = safe_get(target_attachment, 'content')
            if not content:
                return AttachmentContent(
                    attachment=attachment,
                    success=False,
                    error="No content available"
                )

            # Parse based on type
            if attachment.attachment_type == AttachmentType.DOCUMENT:
                if attachment.extension == '.pdf':
                    return await self._parse_pdf(attachment, content, extract_text, extract_images, max_pages)
                elif attachment.extension in ['.docx', '.doc']:
                    return await self._parse_docx(attachment, content, extract_text, extract_images)
                elif attachment.extension in ['.txt', '.rtf']:
                    return await self._parse_text(attachment, content)

            elif attachment.attachment_type == AttachmentType.SPREADSHEET:
                if attachment.extension in ['.xlsx', '.xls']:
                    return await self._parse_excel(attachment, content)
                elif attachment.extension == '.csv':
                    return await self._parse_csv(attachment, content)

            elif attachment.attachment_type == AttachmentType.PRESENTATION:
                if attachment.extension in ['.pptx', '.ppt']:
                    return await self._parse_pptx(attachment, content, extract_text, extract_images)

            elif attachment.attachment_type == AttachmentType.ARCHIVE:
                if attachment.extension == '.zip':
                    return await self._parse_zip(attachment, content)

            # Default: return raw content info
            return AttachmentContent(
                attachment=attachment,
                format=attachment.extension or 'unknown',
                success=True,
                warnings=[f"Parser not available for {attachment.extension}"]
            )

        except Exception as e:
            self.logger.error(f"Failed to read attachment: {e}")
            return None

    async def _get_message(self, message_id: str):
        """Get message from any folder."""
        for folder_name in ['inbox', 'sent', 'drafts']:
            folder = getattr(self.ews_client.account, folder_name, None)
            if folder is None:
                continue
            try:
                message = folder.get(id=message_id)
                if message:
                    return message
            except Exception as e:
                self.logger.debug(f"Message not in {folder_name}: {e}")
                continue
        return None

    async def _parse_pdf(
        self,
        attachment: Attachment,
        content: bytes,
        extract_text: bool,
        extract_images: bool,
        max_pages: int
    ) -> AttachmentContent:
        """Parse PDF file."""
        try:
            import pdfplumber

            result = AttachmentContent(attachment=attachment, format='pdf')

            with pdfplumber.open(io.BytesIO(content)) as pdf:
                # Extract text
                if extract_text:
                    pages_to_process = min(len(pdf.pages), max_pages)
                    text_parts = []
                    for i in range(pages_to_process):
                        page_text = pdf.pages[i].extract_text() or ''
                        text_parts.append(page_text)

                    result.text = '\n\n'.join(text_parts)

                # Extract tables
                result.tables = []
                for i in range(min(len(pdf.pages), max_pages)):
                    tables = pdf.pages[i].extract_tables()
                    if tables:
                        result.tables.extend(tables)

                # Metadata
                if pdf.metadata:
                    result.metadata = {
                        'pages': len(pdf.pages),
                        'author': pdf.metadata.get('Author'),
                        'title': pdf.metadata.get('Title'),
                        'created': str(pdf.metadata.get('CreationDate')),
                    }

            result.success = True
            return result

        except ImportError:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error="pdfplumber not installed. Install with: pip install pdfplumber"
            )
        except Exception as e:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error=f"Failed to parse PDF: {str(e)}"
            )

    async def _parse_docx(
        self,
        attachment: Attachment,
        content: bytes,
        extract_text: bool,
        extract_images: bool
    ) -> AttachmentContent:
        """Parse DOCX file."""
        try:
            from docx import Document

            result = AttachmentContent(attachment=attachment, format='docx')

            doc = Document(io.BytesIO(content))

            # Extract text
            if extract_text:
                paragraphs = [para.text for para in doc.paragraphs]
                result.text = '\n\n'.join(paragraphs)

            # Extract tables
            result.tables = []
            for table in doc.tables:
                table_data = [
                    [cell.text for cell in row.cells]
                    for row in table.rows
                ]
                result.tables.append(table_data)

            # Metadata
            result.metadata = {
                'author': doc.core_properties.author,
                'title': doc.core_properties.title,
                'created': str(doc.core_properties.created),
                'modified': str(doc.core_properties.modified),
            }

            result.success = True
            return result

        except ImportError:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error="python-docx not installed. Install with: pip install python-docx"
            )
        except Exception as e:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error=f"Failed to parse DOCX: {str(e)}"
            )

    async def _parse_excel(
        self,
        attachment: Attachment,
        content: bytes
    ) -> AttachmentContent:
        """Parse Excel file."""
        try:
            from openpyxl import load_workbook

            result = AttachmentContent(attachment=attachment, format='excel')

            wb = load_workbook(io.BytesIO(content), data_only=True)

            # Extract all sheets
            result.sheets = {}
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    sheet_data.append(list(row))
                result.sheets[sheet_name] = sheet_data

            # Try to extract first sheet as structured data
            if wb.sheetnames:
                first_sheet = wb[wb.sheetnames[0]]
                headers = [cell.value for cell in first_sheet[1]]

                result.structured_data = []
                for row in first_sheet.iter_rows(min_row=2, values_only=True):
                    row_dict = dict(zip(headers, row))
                    result.structured_data.append(row_dict)

            result.success = True
            return result

        except ImportError:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error="openpyxl not installed. Install with: pip install openpyxl"
            )
        except Exception as e:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error=f"Failed to parse Excel: {str(e)}"
            )

    async def _parse_text(
        self,
        attachment: Attachment,
        content: bytes
    ) -> AttachmentContent:
        """Parse plain text file."""
        try:
            text = content.decode('utf-8')
            return AttachmentContent(
                attachment=attachment,
                format='text',
                text=text,
                success=True
            )
        except Exception as e:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error=f"Failed to parse text: {str(e)}"
            )

    async def _parse_csv(
        self,
        attachment: Attachment,
        content: bytes
    ) -> AttachmentContent:
        """Parse CSV file."""
        try:
            import csv

            result = AttachmentContent(attachment=attachment, format='csv')

            text = content.decode('utf-8')
            reader = csv.DictReader(io.StringIO(text))

            result.structured_data = list(reader)
            result.success = True
            return result

        except Exception as e:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error=f"Failed to parse CSV: {str(e)}"
            )

    async def _parse_pptx(
        self,
        attachment: Attachment,
        content: bytes,
        extract_text: bool,
        extract_images: bool
    ) -> AttachmentContent:
        """Parse PPTX file."""
        try:
            from pptx import Presentation

            result = AttachmentContent(attachment=attachment, format='pptx')

            prs = Presentation(io.BytesIO(content))

            # Extract text from slides
            if extract_text:
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text = f"--- Slide {i+1} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text += shape.text + "\n"
                    slides_text.append(slide_text)

                result.text = '\n\n'.join(slides_text)

            # Metadata
            result.metadata = {
                'slides': len(prs.slides),
                'author': prs.core_properties.author,
                'title': prs.core_properties.title,
            }

            result.success = True
            return result

        except ImportError:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error="python-pptx not installed. Install with: pip install python-pptx"
            )
        except Exception as e:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error=f"Failed to parse PPTX: {str(e)}"
            )

    async def _parse_zip(
        self,
        attachment: Attachment,
        content: bytes
    ) -> AttachmentContent:
        """Parse ZIP archive."""
        try:
            import zipfile

            result = AttachmentContent(attachment=attachment, format='zip')

            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                # List files
                result.files = []
                for info in zf.infolist():
                    result.files.append({
                        'name': info.filename,
                        'size': info.file_size,
                        'compressed_size': info.compress_size,
                        'modified': str(info.date_time),
                    })

                # Extract readable text files
                result.extracted_content = {}
                for file_info in zf.infolist():
                    if file_info.filename.endswith(('.txt', '.md', '.csv', '.json')):
                        try:
                            with zf.open(file_info) as f:
                                content_text = f.read().decode('utf-8')
                                result.extracted_content[file_info.filename] = content_text
                        except Exception as e:
                            self.logger.debug(f"Skipped {file_info.filename}: {e}")

            result.success = True
            return result

        except Exception as e:
            return AttachmentContent(
                attachment=attachment,
                success=False,
                error=f"Failed to parse ZIP: {str(e)}"
            )
